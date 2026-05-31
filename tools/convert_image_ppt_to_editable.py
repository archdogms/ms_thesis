# -*- coding: utf-8 -*-
"""Convert the image-only Scheme 2 defense deck into an editable PPTX.

The source image-only deck was created from a PIL drawing script.  Instead of
trying to OCR every rendered pixel from scratch, this converter replays that
same drawing script into python-pptx primitives: text boxes, editable shapes,
lines/freeforms, and only local raster assets for complex charts/maps.
"""
from __future__ import annotations

import importlib.util
import json
import math
import os
import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


ROOT = Path(r"C:\Users\31484\Desktop\thesis")
PPT_DIR = ROOT / "ms_thesis" / "docs" / "ppt"
INPUT_PPTX = PPT_DIR / "孟帅_本科毕业设计答辩_浅蓝绿灰方案二升级版_内容图禁紫_整页图片版.pptx"
OUTPUT_PPTX = PPT_DIR / "editable_output.pptx"
REPORT_MD = PPT_DIR / "report.md"
EXTRACTED_DIR = PPT_DIR / "extracted_slides"
SOURCE_SCRIPT = ROOT / "outputs" / "manual-20260526-defense" / "scripts" / "render_scheme2_bluegreen_slides.py"

W, H = 1280, 720
PX_EMU = 914400 / 96
PX_PT = 72 / 96
FONT_CJK = "Microsoft YaHei"
FONT_LATIN = "Arial"


def px(value: float | int) -> Emu:
    return Emu(int(round(float(value) * PX_EMU)))


def pt_from_px(value: float | int) -> Pt:
    return Pt(float(value) * PX_PT)


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.strip().lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def _blend_rgba(r: int, g: int, b: int, a: int, base=(255, 255, 255)) -> tuple[int, int, int]:
    alpha = max(0, min(255, a)) / 255.0
    return tuple(int(round(c * alpha + bc * (1 - alpha))) for c, bc in zip((r, g, b), base))


def color_to_rgb(value: Any, default: str = "#000000") -> tuple[int, int, int]:
    if value is None:
        return _hex_to_rgb(default)
    if isinstance(value, str):
        if value.lower() in {"white", "black"}:
            return (255, 255, 255) if value.lower() == "white" else (0, 0, 0)
        return _hex_to_rgb(value)
    if isinstance(value, (tuple, list)):
        if len(value) >= 4:
            return _blend_rgba(int(value[0]), int(value[1]), int(value[2]), int(value[3]))
        if len(value) >= 3:
            return tuple(int(v) for v in value[:3])
    return _hex_to_rgb(default)


def set_solid_fill(shape, fill: Any) -> None:
    if fill is None:
        shape.fill.background()
        return
    r, g, b = color_to_rgb(fill, "#FFFFFF")
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def set_line(shape, outline: Any, width: float | int = 1) -> None:
    if outline is None or width == 0:
        shape.line.fill.background()
        shape.line.width = Pt(0)
        return
    r, g, b = color_to_rgb(outline, "#000000")
    shape.line.color.rgb = RGBColor(r, g, b)
    shape.line.width = pt_from_px(width)


def set_text_font(run, size_px: int, color: Any, bold: bool) -> None:
    font = run.font
    font.name = FONT_CJK
    font.size = pt_from_px(size_px)
    font.bold = bool(bold)
    r, g, b = color_to_rgb(color, "#2B2F33")
    font.color.rgb = RGBColor(r, g, b)

    rpr = run._r.get_or_add_rPr()
    rfonts = rpr.find(qn("a:rFonts"))
    if rfonts is None:
        rfonts = OxmlElement("a:rFonts")
        rpr.insert(0, rfonts)
    rfonts.set("latin", FONT_LATIN)
    rfonts.set("ea", FONT_CJK)
    rfonts.set("cs", FONT_CJK)


@dataclass
class SlideStats:
    text_boxes: int = 0
    shapes: int = 0
    lines: int = 0
    pictures: int = 0
    freeforms: int = 0
    raster_assets: list[str] = field(default_factory=list)
    strategy: str = ""


class EditableImage:
    def __init__(self, deck: "EditableDeck"):
        self.deck = deck

    def paste(self, *_args, **_kwargs):
        # All semantic paste operations are intercepted by paste_logo/fit_paste.
        return None

    def save(self, path: str | os.PathLike[str], *args, **kwargs):
        # Debug convenience: save a blank placeholder if any caller expects save().
        Image.new("RGB", (W, H), self.deck.mod.COL["bg"]).save(path, *args, **kwargs)


class PptxDraw:
    def __init__(self, deck: "EditableDeck"):
        self.deck = deck

    @property
    def slide(self):
        return self.deck.slide

    @property
    def stats(self) -> SlideStats:
        return self.deck.current_stats

    def rectangle(self, box, fill=None, outline=None, width=1):
        x1, y1, x2, y2 = [float(v) for v in box]
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x1), px(y1), px(x2 - x1), px(y2 - y1))
        set_solid_fill(shape, fill)
        set_line(shape, outline, width)
        self.stats.shapes += 1
        return shape

    def rounded_rectangle(self, box, radius=16, fill=None, outline=None, width=1):
        x1, y1, x2, y2 = [float(v) for v in box]
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x1), px(y1), px(x2 - x1), px(y2 - y1))
        set_solid_fill(shape, fill)
        set_line(shape, outline, width)
        self.stats.shapes += 1
        return shape

    def ellipse(self, box, fill=None, outline=None, width=1):
        x1, y1, x2, y2 = [float(v) for v in box]
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x1), px(y1), px(x2 - x1), px(y2 - y1))
        set_solid_fill(shape, fill)
        set_line(shape, outline, width)
        self.stats.shapes += 1
        return shape

    def polygon(self, points, fill=None, outline=None, width=1):
        pts = [(float(x), float(y)) for x, y in points]
        if len(pts) < 3:
            return None
        shape = self._freeform(pts, close=True)
        set_solid_fill(shape, fill)
        set_line(shape, outline, width)
        self.stats.shapes += 1
        self.stats.freeforms += 1
        return shape

    def line(self, xy, fill=None, width=1, joint=None):
        pts = normalize_line_points(xy)
        if len(pts) < 2:
            return None
        if len(pts) == 2:
            x1, y1 = pts[0]
            x2, y2 = pts[1]
            line = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
            set_line(line, fill, width)
            self.stats.lines += 1
            return line
        shape = self._freeform(pts, close=False)
        shape.fill.background()
        set_line(shape, fill, width)
        self.stats.lines += 1
        self.stats.freeforms += 1
        return shape

    def arc(self, box, start, end, fill=None, width=1):
        x1, y1, x2, y2 = [float(v) for v in box]
        # PowerPoint arc control is limited through python-pptx; use a transparent
        # block arc as a close editable approximation for decorative cover rings.
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ARC, px(x1), px(y1), px(x2 - x1), px(y2 - y1))
        shape.fill.background()
        set_line(shape, fill, width)
        self.stats.shapes += 1
        return shape

    def multiline_text(self, xy, text, font=None, fill=None, spacing=0, anchor=None, align="left"):
        size = getattr(font, "size", 18) if font is not None else 18
        self.deck.add_text(xy, text, size_px=size, color=fill, bold=False, max_width=None, anchor=anchor, align=align)

    def textbbox(self, xy, text, font=None):
        # Only used by the original script's measure helper when it receives our
        # draw object by accident; return a conservative text box.
        size = getattr(font, "size", 18) if font is not None else 18
        return (0, 0, int(len(str(text)) * size * 0.55), int(size * 1.2))

    def _freeform(self, pts: list[tuple[float, float]], close: bool):
        start_x, start_y = pts[0]
        builder = self.slide.shapes.build_freeform(px(start_x), px(start_y))
        builder.add_line_segments([(px(x), px(y)) for x, y in pts[1:]], close=close)
        return builder.convert_to_shape()


def normalize_line_points(xy) -> list[tuple[float, float]]:
    if isinstance(xy, (tuple, list)) and len(xy) == 4 and all(isinstance(v, (int, float)) for v in xy):
        return [(float(xy[0]), float(xy[1])), (float(xy[2]), float(xy[3]))]
    return [(float(p[0]), float(p[1])) for p in xy]


class EditableDeck:
    def __init__(self, mod):
        self.mod = mod
        self.prs = Presentation()
        self.prs.slide_width = px(W)
        self.prs.slide_height = px(H)
        self.slide = None
        self.current_stats = SlideStats()
        self.stats: list[SlideStats] = []

    def new_slide(self):
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_stats = SlideStats()
        draw = PptxDraw(self)

        bg = self.slide.background.fill
        bg.solid()
        r, g, b = color_to_rgb(self.mod.COL["bg"], "#F6F9FA")
        bg.fore_color.rgb = RGBColor(r, g, b)

        # Quiet cartographic grid from the source renderer.
        for x in range(0, W, 80):
            draw.line((x, 0, x, H), fill="#E8EEF1", width=1)
        for y in range(0, H, 80):
            draw.line((0, y, W, y), fill="#E8EEF1", width=1)
        return EditableImage(self), draw

    def add_text(
        self,
        xy,
        text,
        size_px=24,
        color=None,
        bold=False,
        max_width=None,
        line_spacing=1.18,
        anchor=None,
        align="left",
    ):
        text = str(text)
        font = self.mod.f(size_px, bold=bold)
        value = self.mod.wrap_text(text, font, max_width) if max_width else text
        lines = value.split("\n") or [""]

        measured_width = max((self.mod.measure(line, font)[0] for line in lines), default=size_px)
        width = float(max_width or measured_width + 8)
        height = float(max(size_px * line_spacing * len(lines) + 8, size_px + 8))
        x, y = float(xy[0]), float(xy[1])

        if anchor:
            if anchor[0] == "m":
                x -= width / 2
            elif anchor[0] == "r":
                x -= width
            if len(anchor) > 1 and anchor[1] == "m":
                y -= height / 2

        tb = self.slide.shapes.add_textbox(px(x), px(y), px(width), px(height))
        tb.name = f"editable-text-{self.current_stats.text_boxes + 1:03d}"
        tf = tb.text_frame
        tf.clear()
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = True
        if anchor and len(anchor) > 1 and anchor[1] == "m":
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = {
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "left": PP_ALIGN.LEFT,
            }.get(align, PP_ALIGN.LEFT)
            p.line_spacing = line_spacing
            run = p.add_run()
            run.text = line
            set_text_font(run, size_px, color or self.mod.COL["ink"], bold)

        self.current_stats.text_boxes += 1
        return tb

    def add_picture_fit(self, path: Path, box, bg="#FFFFFF", border=True):
        x, y, w, h = [float(v) for v in box]
        draw = PptxDraw(self)
        draw.rounded_rectangle((x, y, x + w, y + h), radius=12, fill=bg, outline=self.mod.COL["line"] if border else bg, width=1)

        im = Image.open(path)
        iw, ih = im.size
        scale = min((w - 28) / iw, (h - 28) / ih)
        rw, rh = iw * scale, ih * scale
        px_pos = x + (w - rw) / 2
        py_pos = y + (h - rh) / 2
        pic = self.slide.shapes.add_picture(str(path), px(px_pos), px(py_pos), width=px(rw), height=px(rh))
        pic.name = f"raster-asset-{self.current_stats.pictures + 1:03d}-{path.name}"
        self.current_stats.pictures += 1
        self.current_stats.raster_assets.append(path.name)
        return pic

    def add_logo(self, x=1046, y=34, w=170):
        logo_path = self.mod.ASSET / "tsinghua-logo.png"
        im = Image.open(logo_path)
        ratio = im.height / im.width
        pic = self.slide.shapes.add_picture(str(logo_path), px(x), px(y), width=px(w), height=px(w * ratio))
        pic.name = "raster-logo-tsinghua"
        self.current_stats.pictures += 1
        self.current_stats.raster_assets.append("tsinghua-logo.png")
        return pic

    def finalize_slide(self, strategy: str):
        self.current_stats.strategy = strategy
        self.stats.append(self.current_stats)


def load_source_module():
    spec = importlib.util.spec_from_file_location("scheme2_bluegreen_renderer", SOURCE_SCRIPT)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot import source script: {SOURCE_SCRIPT}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def inspect_and_extract_input():
    EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
    for old in EXTRACTED_DIR.glob("*"):
        if old.is_file():
            old.unlink()

    src = Presentation(INPUT_PPTX)
    info = {
        "input": str(INPUT_PPTX),
        "slide_count": len(src.slides),
        "slide_width_emu": src.slide_width,
        "slide_height_emu": src.slide_height,
        "slide_width_in": round(src.slide_width / 914400, 4),
        "slide_height_in": round(src.slide_height / 914400, 4),
        "slides": [],
    }

    for idx, slide in enumerate(src.slides, 1):
        slide_info = {"index": idx, "shape_count": len(slide.shapes), "shapes": [], "extracted": None}
        for sh in slide.shapes:
            rec = {
                "type": str(sh.shape_type),
                "name": getattr(sh, "name", ""),
                "left": int(sh.left),
                "top": int(sh.top),
                "width": int(sh.width),
                "height": int(sh.height),
            }
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                rec["is_picture"] = True
                rec["coverage"] = round((sh.width * sh.height) / (src.slide_width * src.slide_height), 4)
                if rec["coverage"] >= 0.95:
                    ext = sh.image.ext or "png"
                    out = EXTRACTED_DIR / f"slide_{idx:02d}.{ext}"
                    out.write_bytes(sh.image.blob)
                    slide_info["extracted"] = str(out)
            slide_info["shapes"].append(rec)
        info["slides"].append(slide_info)
    return info


def build_editable_deck():
    mod = load_source_module()
    deck = EditableDeck(mod)

    def new_slide_override():
        return deck.new_slide()

    def paste_logo_override(img, x=1046, y=34, w=170):
        return deck.add_logo(x, y, w)

    def fit_paste_override(img, path, box, bg="#FFFFFF", border=True):
        return deck.add_picture_fit(Path(path), box, bg, border)

    def draw_text_override(d, xy, text, size=24, color=None, bold=False, max_width=None, line_spacing=1.18, anchor=None, align="left"):
        return deck.add_text(xy, text, size, color or mod.COL["ink"], bold, max_width, line_spacing, anchor, align)

    mod.new_slide = new_slide_override
    mod.paste_logo = paste_logo_override
    mod.fit_paste = fit_paste_override
    mod.draw_text = draw_text_override

    strategies = {
        1: "Cover rebuilt with editable title, metadata, label, geometric panels, grid, and editable Nanhai outline; Tsinghua logo retained as raster identity asset.",
        2: "Timeline rebuilt with editable line, nodes, labels, and bottom statement.",
        3: "Two-sided concept comparison rebuilt with editable panels, rules, central rounded rectangle, and all text.",
        4: "Object slide rebuilt with editable Nanhai outline, metrics, bar chart primitives, and note box.",
        5: "Research questions rebuilt as editable radial node diagram with lines, circles, labels, and conclusion text.",
        6: "Data foundation rebuilt with editable metrics, connectors, center node, output labels, and arrow.",
        7: "Technical route rebuilt with editable input/output boxes, route nodes, connectors, and text; lower pipeline illustration retained as local raster asset.",
        8: "Culture-side method rebuilt with editable metric rail, footer proof text, and local knowledge-graph example image.",
        9: "Tourism-side status rebuilt with editable KPI block and explanatory text; internal POI charts retained as local raster assets.",
        10: "Bridge and indicators rebuilt with editable circles, bridge box, metric cards, labels, and connector line.",
        11: "Quadrant matrix rebuilt with editable rectangles, axis labels, quadrant counts, side callout, and stacked bar primitives.",
        12: "Spatial result rebuilt with editable side explanation and planning label; density map retained as local raster asset.",
        13: "Grid validation rebuilt with editable comparison card, arrow, rows, values, and conclusion label; map comparison retained as local raster asset.",
        14: "Jiujiang case rebuilt with editable explanation panel and metrics; zoom map retained as local raster asset.",
        15: "Robustness slide rebuilt with editable central badge, coefficients, and response banner; two analytical charts retained as local raster assets.",
        16: "Comment review rebuilt with editable explanation panel, metrics, and boundary label; comment chart retained as local raster asset.",
        17: "Planning translation rebuilt as editable three-lane flow diagram with round rectangles, arrows, and labels.",
        18: "Closing slide rebuilt with editable two-column panels, numbered innovation list, boundary list, conclusion box, and thanks text.",
    }

    for idx, builder in enumerate(mod.SLIDES, 1):
        builder()
        deck.finalize_slide(strategies[idx])

    deck.prs.save(OUTPUT_PPTX)
    return deck.stats


def validate_output():
    prs = Presentation(OUTPUT_PPTX)
    validation = {
        "output": str(OUTPUT_PPTX),
        "exists": OUTPUT_PPTX.exists(),
        "bytes": OUTPUT_PPTX.stat().st_size if OUTPUT_PPTX.exists() else 0,
        "slide_count": len(prs.slides),
        "slides": [],
    }
    for idx, slide in enumerate(prs.slides, 1):
        text_shapes = 0
        picture_shapes = 0
        shape_count = len(slide.shapes)
        full_slide_pictures = 0
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text.strip():
                text_shapes += 1
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if sh.left <= px(2) and sh.top <= px(2) and sh.width >= px(W - 4) and sh.height >= px(H - 4):
                    full_slide_pictures += 1
        validation["slides"].append(
            {
                "index": idx,
                "shape_count": shape_count,
                "text_shapes": text_shapes,
                "pictures": picture_shapes,
                "full_slide_pictures": full_slide_pictures,
            }
        )

    with zipfile.ZipFile(OUTPUT_PPTX) as zf:
        media = [i for i in zf.infolist() if i.filename.startswith("ppt/media/")]
        validation["media_count"] = len(media)
        validation["empty_media"] = [m.filename for m in media if m.file_size == 0]
    return validation


def write_report(input_info, stats: list[SlideStats], validation):
    lines = []
    lines.append("# 图片型 PPT 转可编辑 PPT 报告")
    lines.append("")
    lines.append(f"- 输入文件: `{INPUT_PPTX}`")
    lines.append(f"- 输出文件: `{OUTPUT_PPTX}`")
    lines.append(f"- 提取图片目录: `{EXTRACTED_DIR}`")
    lines.append(f"- 页面尺寸: {input_info['slide_width_in']} x {input_info['slide_height_in']} in")
    lines.append(f"- 页数: {input_info['slide_count']}")
    lines.append("")
    lines.append("## 输入结构检查")
    lines.append("")
    for slide in input_info["slides"]:
        full = "是" if slide["extracted"] else "否"
        lines.append(f"- Slide {slide['index']:02d}: shape_count={slide['shape_count']}, 整页图片={full}, extracted={slide['extracted'] or '无'}")
    lines.append("")
    lines.append("## 重建策略")
    lines.append("")
    for idx, st in enumerate(stats, 1):
        raster = ", ".join(sorted(set(st.raster_assets))) if st.raster_assets else "无"
        lines.append(f"### Slide {idx:02d}")
        lines.append("")
        lines.append(f"- 策略: {st.strategy}")
        lines.append(f"- 可编辑文本框: {st.text_boxes}")
        lines.append(f"- 可编辑 shape: {st.shapes}")
        lines.append(f"- 可编辑线条/自由线: {st.lines}")
        lines.append(f"- 局部保留图片: {st.pictures} ({raster})")
        lines.append("")
    lines.append("## 输出校验")
    lines.append("")
    lines.append(f"- python-pptx 可正常读取: {'是' if validation['exists'] and validation['slide_count'] == input_info['slide_count'] else '否'}")
    lines.append(f"- 输出大小: {validation['bytes']} bytes")
    lines.append(f"- 输出页数: {validation['slide_count']}")
    lines.append(f"- 媒体文件数: {validation['media_count']}")
    lines.append(f"- 空媒体文件: {validation['empty_media'] or '无'}")
    lines.append("")
    lines.append("| 页码 | 总对象 | 文本框 | 图片 | 整页图片 |")
    lines.append("|---:|---:|---:|---:|---:|")
    for s in validation["slides"]:
        lines.append(f"| {s['index']:02d} | {s['shape_count']} | {s['text_shapes']} | {s['pictures']} | {s['full_slide_pictures']} |")
    lines.append("")
    lines.append("## 说明")
    lines.append("")
    lines.append("- 新文件不是把原始整页图片贴回去；校验中每页 `整页图片` 均为 0。")
    lines.append("- 标题、正文、数字、标签、KPI、流程文字、解释文字均由 PowerPoint 文本框重建，可编辑。")
    lines.append("- 背景网格、色块、矩形、圆角矩形、圆形、线条、箭头、矩阵和流程框均由 PowerPoint shape/freeform 重建。")
    lines.append("- 复杂地图/图表/知识图谱示意等局部视觉对象保留为图片，并在每页策略中列明。")
    REPORT_MD.write_text("\n".join(lines), encoding="utf-8")


def main():
    if not INPUT_PPTX.exists():
        raise FileNotFoundError(INPUT_PPTX)
    if not SOURCE_SCRIPT.exists():
        raise FileNotFoundError(SOURCE_SCRIPT)
    input_info = inspect_and_extract_input()
    stats = build_editable_deck()
    validation = validate_output()
    write_report(input_info, stats, validation)
    print(json.dumps({"output": str(OUTPUT_PPTX), "report": str(REPORT_MD), "validation": validation}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
